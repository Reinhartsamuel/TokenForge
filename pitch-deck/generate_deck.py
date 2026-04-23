from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

SOL_GREEN = RGBColor(0x14, 0xF1, 0x95)
SOL_PURPLE = RGBColor(0x99, 0x45, 0xFF)
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)
DARK_CARD = RGBColor(0x1A, 0x1A, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0x8B, 0x8B, 0xA7)
MED_GRAY = RGBColor(0x4A, 0x4A, 0x6A)
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xFF)
RED_ACCENT = RGBColor(0xFF, 0x4D, 0x4D)
ORANGE_ACCENT = RGBColor(0xFF, 0x8C, 0x00)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color if fill_color else DARK_CARD
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.adjustments[0] = 0.05
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=18, color=WHITE, bold=False, line_spacing=1.3, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.4)
        p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=None, spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(font_size * 0.6)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.findall(qn("a:buNone"))
        for bn in buNone:
            pPr.remove(bn)
        buChar = pPr.find(qn("a:buChar"))
        if buChar is None:
            from lxml import etree
            buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "\u2022")
        if bullet_color:
            buClr = pPr.find(qn("a:buClr"))
            if buClr is None:
                from lxml import etree
                buClr = etree.SubElement(pPr, qn("a:buClr"))
            else:
                for child in list(buClr):
                    buClr.remove(child)
            from lxml import etree
            srgbClr = etree.SubElement(buClr, qn("a:srgbClr"))
            srgbClr.set("val", "{:02X}{:02X}{:02X}".format(bullet_color[0], bullet_color[1], bullet_color[2]))
    return txBox


def add_accent_line(slide, left, top, width, color=SOL_GREEN, height=Pt(3)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gradient_bar(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = SOL_GREEN
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_accent_line(slide, Inches(1.2), Inches(2.2), Inches(0.8), SOL_GREEN, Pt(4))

add_text(slide, Inches(1.2), Inches(2.4), Inches(8), Inches(1.2), "TokenForge", 54, WHITE, True)
add_text(slide, Inches(1.2), Inches(3.5), Inches(8), Inches(0.8), "The Metaplex of RWA Tokens", 28, SOL_GREEN, False)

add_multiline(slide, Inches(1.2), Inches(4.5), Inches(8), Inches(1.5), [
    "The first open-source SDK and issuer dashboard for the",
    "Solana Security Token Standard (SSTS)"
], 18, LIGHT_GRAY, False)

add_text(slide, Inches(1.2), Inches(6.2), Inches(4), Inches(0.4), "April 2026  |  Confidential", 12, MED_GRAY)

# Purple accent circle top-right
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.5), Inches(0.5), Inches(2.5), Inches(2.5))
shape.fill.solid()
shape.fill.fore_color.rgb = SOL_PURPLE
shape.line.fill.background()
# Make semi-transparent via alpha
from lxml import etree
from pptx.oxml.ns import qn
spPr = shape._element.spPr
solidFill_el = spPr.find(qn("a:solidFill"))
if solidFill_el is not None:
    srgbClr = solidFill_el.find(qn("a:srgbClr"))
    if srgbClr is not None:
        alpha = etree.SubElement(srgbClr, qn("a:alpha"))
        alpha.set("val", "15000")

shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.8), Inches(4.5), Inches(3), Inches(3))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SOL_GREEN
shape2.line.fill.background()
spPr2 = shape2._element.spPr
solidFill_el2 = spPr2.find(qn("a:solidFill"))
if solidFill_el2 is not None:
    srgbClr2 = solidFill_el2.find(qn("a:srgbClr"))
    if srgbClr2 is not None:
        alpha2 = etree.SubElement(srgbClr2, qn("a:alpha"))
        alpha2.set("val", "10000")


# ═══════════════════════════════════════════════════════
# SLIDE 2: PROBLEM
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "THE PROBLEM", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
         "Issuing a Compliant Security Token\non Solana Takes 2\u20133 Months", 38, WHITE, True)

add_bullet_list(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.5), [
    "Halborn SSTS spec published Jan 2026 \u2014 zero tooling exists three months later",
    "Asset managers must hand-roll Anchor CPI calls, implement KYC/AML verification, build Merkle trees for distributions, and stand up admin UIs",
    "Franklin Templeton, BlackRock, Securitize are all live on Solana \u2014 but there\u2019s no SDK to follow",
    "2\u20133 month engineering project before a single asset goes live"
], 16, LIGHT_GRAY, (0x14, 0xF1, 0x95))

# Right side: cost card
add_rect(slide, Inches(7.2), Inches(3.0), Inches(5.2), Inches(3.5), DARK_CARD, SOL_GREEN, Pt(1))
add_text(slide, Inches(7.6), Inches(3.3), Inches(4.4), Inches(0.5), "Cost of Manual Implementation", 16, SOL_GREEN, True)
add_text(slide, Inches(7.6), Inches(3.9), Inches(4.4), Inches(1.0), "$50K+", 60, WHITE, True)
add_multiline(slide, Inches(7.6), Inches(5.0), Inches(4.4), Inches(1.5), [
    "2\u20133 months of engineering time",
    "No dashboard, no SDK, no canonical FAMP",
    "Repeated by every issuer independently"
], 14, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 3: AGITATE - SSTS EXISTS BUT NO TOOLING
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "THE GAP", 12, RED_ACCENT, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), RED_ACCENT, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.2),
         "The SSTS Spec Exists.\nThe Tooling Doesn\u2019t.", 38, WHITE, True)

# Left column: What exists
add_rect(slide, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.8), DARK_CARD)
add_text(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(0.5), "What Exists", 18, SOL_GREEN, True)
add_multiline(slide, Inches(1.2), Inches(3.8), Inches(4.8), Inches(2.8), [
    "\u2713  SSTS spec by Halborn (Jan 2026)",
    "\u2713  Co-authored by Franklin Templeton,",
    "    Securitize, Superstate, Solana Foundation",
    "\u2713  FAMP proposal (sRFC 37)",
    "\u2713  Major institutions live on Solana"
], 16, LIGHT_GRAY)

# Right column: What's missing
add_rect(slide, Inches(6.8), Inches(3.0), Inches(5.5), Inches(3.8), DARK_CARD, RED_ACCENT, Pt(1))
add_text(slide, Inches(7.2), Inches(3.2), Inches(4.8), Inches(0.5), "What\u2019s Missing", 18, RED_ACCENT, True)
add_multiline(slide, Inches(7.2), Inches(3.8), Inches(4.8), Inches(2.8), [
    "\u2717  No SSTS implementation program",
    "\u2717  No FAMP canonical program",
    "\u2717  No TypeScript SDK",
    "\u2717  No issuer dashboard",
    "\u2717  Zero builder tooling exists"
], 16, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 4: SOLUTION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "THE SOLUTION", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(1.0),
         "From Months to One Day", 38, WHITE, True)

add_multiline(slide, Inches(0.8), Inches(2.5), Inches(11), Inches(1.0), [
    "TokenForge wraps the full SSTS Anchor program + FAMP into a typed TypeScript SDK",
    "and a no-code issuer dashboard \u2014 compliance teams can mint, manage, and distribute without writing Rust."
], 18, LIGHT_GRAY)

# Three feature cards
card_data = [
    ("SSTS Anchor Program", "Full SSTS spec implementation: create_token, mint_to, transfer hook, freeze/thaw, pause, distributions, corporate actions", SOL_GREEN),
    ("FAMP Canonical Impl.", "sRFC 37 Freeze Authority Management: block/allow lists as Merkle roots, atomic enforcement, delegated freeze authority", SOL_PURPLE),
    ("TypeScript SDK + Dashboard", "npm install @tokenforge/sdk \u2014 typed client for all SSTS + FAMP instructions. No-code dashboard for compliance officers.", ACCENT_BLUE),
]

for i, (title, desc, accent) in enumerate(card_data):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(4.0), Inches(3.7), Inches(3.0), DARK_CARD, accent, Pt(1))
    add_accent_line(slide, Emu(x + Inches(0.3)), Inches(4.3), Inches(0.6), accent, Pt(3))
    add_text(slide, Emu(x + Inches(0.3)), Inches(4.55), Inches(3.1), Inches(0.5), title, 16, accent, True)
    add_text(slide, Emu(x + Inches(0.3)), Inches(5.1), Inches(3.1), Inches(1.8), desc, 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 5: HOW IT WORKS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "HOW IT WORKS", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Three Steps to Compliant Token Issuance", 38, WHITE, True)

steps = [
    ("1", "Connect + Configure", "Issuer connects a multisig wallet, sets the KYC verification program address and FAMP block/allow list policy. TokenForge deploys the SSTS token with all extensions pre-configured.", SOL_GREEN),
    ("2", "Mint + Manage", "Mint tokens to KYC\u2019d holders from the dashboard. Freeze/thaw individual accounts or update the block/allow list in one click. All actions are CPI calls to the on-chain SSTS program.", SOL_PURPLE),
    ("3", "Distribute + Report", "Upload a holder snapshot. TokenForge builds the Merkle tree, posts the root on-chain, and holders claim dividends. Dashboard shows holder registry and distribution history.", ACCENT_BLUE),
]

for i, (num, title, desc, accent) in enumerate(steps):
    y = Inches(2.4 + i * 1.65)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)

    add_text(slide, Inches(2.0), y - Inches(0.05), Inches(3), Inches(0.5), title, 20, accent, True)
    add_text(slide, Inches(2.0), y + Inches(0.45), Inches(9.5), Inches(1.0), desc, 15, LIGHT_GRAY)

    # Connecting line between steps
    if i < 2:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Emu(y + Inches(0.75)), Pt(2), Inches(0.9))
        line.fill.solid()
        line.fill.fore_color.rgb = MED_GRAY
        line.line.fill.background()


# ═══════════════════════════════════════════════════════
# SLIDE 6: ARCHITECTURE
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "ARCHITECTURE", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Built for Security and Compliance from Day One", 38, WHITE, True)

# On-chain layer
add_rect(slide, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.8), DARK_CARD, SOL_GREEN, Pt(1))
add_text(slide, Inches(1.2), Inches(2.6), Inches(3), Inches(0.4), "ON-CHAIN  |  Solana Mainnet", 13, SOL_GREEN, True)

add_rect(slide, Inches(1.2), Inches(3.1), Inches(3.5), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(1.4), Inches(3.2), Inches(3.1), Inches(0.6), "SSTS Program\n(Anchor/Rust)", 13, WHITE, True)

add_rect(slide, Inches(5.2), Inches(3.1), Inches(3.5), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(5.4), Inches(3.2), Inches(3.1), Inches(0.6), "FAMP Program\n(sRFC 37)", 13, WHITE, True)

add_rect(slide, Inches(9.2), Inches(3.1), Inches(3), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(9.4), Inches(3.2), Inches(2.6), Inches(0.6), "SPL Token 2022\n(transfer hook)", 13, WHITE, True)

# Off-chain layer
add_rect(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.8), DARK_CARD, SOL_PURPLE, Pt(1))
add_text(slide, Inches(1.2), Inches(4.8), Inches(3), Inches(0.4), "OFF-CHAIN  |  Backend + Edge", 13, SOL_PURPLE, True)

add_rect(slide, Inches(1.2), Inches(5.3), Inches(3.5), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(1.4), Inches(5.4), Inches(3.1), Inches(0.6), "Hono REST API\n(Railway)", 13, WHITE, True)

add_rect(slide, Inches(5.2), Inches(5.3), Inches(3.5), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(5.4), Inches(5.4), Inches(3.1), Inches(0.6), "Merkle Proof Gen\n(Cloudflare Workers)", 13, WHITE, True)

add_rect(slide, Inches(9.2), Inches(5.3), Inches(3), Inches(0.85), RGBColor(0x12, 0x20, 0x30))
add_text(slide, Inches(9.4), Inches(5.4), Inches(2.6), Inches(0.6), "Holder Registry\n(Supabase)", 13, WHITE, True)

# Frontend layer
add_rect(slide, Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.5), DARK_CARD, ACCENT_BLUE, Pt(1))
add_text(slide, Inches(1.2), Inches(6.85), Inches(11), Inches(0.4), "FRONTEND  |  Vite + React + TypeScript  |  Vercel  |  Dashboard + Registry + Claim Pages", 13, ACCENT_BLUE, True)


# ═══════════════════════════════════════════════════════
# SLIDE 7: WHY NOW
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "WHY NOW", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "The Window Is 3\u20139 Months", 38, WHITE, True)

# Timeline visual
add_rect(slide, Inches(1.5), Inches(2.8), Inches(10.5), Inches(0.15), MED_GRAY)
timeline_events = [
    (0.0, "Jan 2026", "SSTS spec\npublished", SOL_GREEN),
    (0.25, "Feb 2026", "Franklin Templeton\n$800M on-chain", SOL_PURPLE),
    (0.45, "Mar 2026", "Contra launches\nzero builder tooling", ACCENT_BLUE),
    (0.7, "Apr 2026", "Today: zero\nSSTS implementations", RED_ACCENT),
    (1.0, "Jul\u2013Dec 2026", "Window closes:\nwell-resourced teams\nclaim this space", ORANGE_ACCENT),
]

for pos, label, desc, color in timeline_events:
    x = Inches(1.5 + pos * 10.5)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(x - Inches(0.12)), Inches(2.65), Inches(0.24), Inches(0.24))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, Emu(x - Inches(0.8)), Inches(3.2), Inches(1.8), Inches(0.3), label, 12, color, True, PP_ALIGN.CENTER)
    add_text(slide, Emu(x - Inches(0.8)), Inches(3.5), Inches(1.8), Inches(1.0), desc, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Key stats
add_rect(slide, Inches(0.8), Inches(5.0), Inches(3.7), Inches(2.0), DARK_CARD, SOL_GREEN, Pt(1))
add_text(slide, Inches(1.1), Inches(5.2), Inches(3.2), Inches(0.4), "RWA Market Growth", 13, SOL_GREEN, True)
add_text(slide, Inches(1.1), Inches(5.6), Inches(3.2), Inches(0.7), "380%", 48, WHITE, True)
add_text(slide, Inches(1.1), Inches(6.3), Inches(3.2), Inches(0.4), "2022\u20132025, reaching $24B", 13, LIGHT_GRAY)

add_rect(slide, Inches(4.9), Inches(5.0), Inches(3.7), Inches(2.0), DARK_CARD, SOL_PURPLE, Pt(1))
add_text(slide, Inches(5.2), Inches(5.2), Inches(3.2), Inches(0.4), "SSTS Co-Authors", 13, SOL_PURPLE, True)
add_text(slide, Inches(5.2), Inches(5.6), Inches(3.2), Inches(0.7), "5", 48, WHITE, True)
add_text(slide, Inches(5.2), Inches(6.3), Inches(3.2), Inches(0.4), "Halborn + Franklin Templeton +\nSecuritize + Superstate + SF", 13, LIGHT_GRAY)

add_rect(slide, Inches(9.0), Inches(5.0), Inches(3.7), Inches(2.0), DARK_CARD, RED_ACCENT, Pt(1))
add_text(slide, Inches(9.3), Inches(5.2), Inches(3.2), Inches(0.4), "Existing Implementations", 13, RED_ACCENT, True)
add_text(slide, Inches(9.3), Inches(5.6), Inches(3.2), Inches(0.7), "0", 48, WHITE, True)
add_text(slide, Inches(9.3), Inches(6.3), Inches(3.2), Inches(0.4), "Zero SSTS tooling exists today", 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 8: MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "MARKET", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "$24B and Growing \u2014 But Zero Tooling", 38, WHITE, True)

# Concentric circles for TAM/SAM/SOM
# TAM
shape_tam = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.5), Inches(4.5), Inches(4.5))
shape_tam.fill.solid()
shape_tam.fill.fore_color.rgb = RGBColor(0x0A, 0x2A, 0x15)
shape_tam.line.color.rgb = SOL_GREEN
shape_tam.line.width = Pt(2)

# SAM
shape_sam = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(3.2), Inches(3.1), Inches(3.1))
shape_sam.fill.solid()
shape_sam.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x35)
shape_sam.line.color.rgb = SOL_PURPLE
shape_sam.line.width = Pt(2)

# SOM
shape_som = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.9), Inches(3.9), Inches(1.7), Inches(1.7))
shape_som.fill.solid()
shape_som.fill.fore_color.rgb = DARK_CARD
shape_som.line.color.rgb = ACCENT_BLUE
shape_som.line.width = Pt(2)

# Labels on circles
add_text(slide, Inches(2.0), Inches(2.6), Inches(2), Inches(0.8), "TAM\n$24B", 14, SOL_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2.5), Inches(3.4), Inches(2), Inches(0.8), "SAM\n$4.8B", 14, SOL_PURPLE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(3.1), Inches(4.3), Inches(1.3), Inches(0.8), "SOM\n$480K", 12, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Right side: breakdown
add_multiline(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(1.0), [
    "TAM: $24B  \u2014  Global RWA tokenization market (2025)"
], 16, SOL_GREEN, True)

add_multiline(slide, Inches(7.0), Inches(3.3), Inches(5.5), Inches(1.0), [
    "SAM: $4.8B  \u2014  Solana RWA segment (20%)"
], 16, SOL_PURPLE, True)

add_multiline(slide, Inches(7.0), Inches(4.1), Inches(5.5), Inches(1.0), [
    "SOM: $480K  \u2014  Year 1 (1% of SAM, 10 issuers)"
], 16, ACCENT_BLUE, True)

add_rect(slide, Inches(7.0), Inches(5.2), Inches(5.5), Inches(1.8), DARK_CARD)
add_text(slide, Inches(7.3), Inches(5.4), Inches(5), Inches(0.4), "Bottom-Up Validation", 14, SOL_GREEN, True)
add_multiline(slide, Inches(7.3), Inches(5.9), Inches(5), Inches(1.0), [
    "10 issuers \u00d7 $500/mo Starter + 2 \u00d7 $2K Growth",
    "= $9K MRR = $108K ARR (Month 3 target)",
    "Break-even at 1\u20132 paying customers"
], 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 9: BUSINESS MODEL
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "BUSINESS MODEL", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Open Source \u2192 SaaS \u2192 Protocol Fee", 38, WHITE, True)

# Three phases
phases = [
    ("Phase 1", "Open Source + Grant", "Month 1", [
        "Full open-source on GitHub",
        "Apply for SF $250K RFP grant",
        "Zero revenue, maximum adoption",
        "Developer mindshare first"
    ], SOL_GREEN),
    ("Phase 2", "Managed Service SaaS", "Month 2\u20136", [
        "Starter: $500/mo (up to $10M AUM)",
        "Growth: $2,000/mo (up to $100M AUM)",
        "Enterprise: $5,000+/mo (unlimited)",
        "Hosted dashboard + Merkle infra"
    ], SOL_PURPLE),
    ("Phase 3", "Protocol Fee", "Month 6+", [
        "Basis-point per-transfer fee",
        "Similar to Metaplex mint fee",
        "Toll road on institutional RWA",
        "Scales with ecosystem growth"
    ], ACCENT_BLUE),
]

for i, (phase, title, timing, items, accent) in enumerate(phases):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(2.5), Inches(3.7), Inches(4.5), DARK_CARD, accent, Pt(1))
    add_text(slide, Emu(x + Inches(0.3)), Inches(2.7), Inches(3.1), Inches(0.3), phase, 11, accent, True)
    add_text(slide, Emu(x + Inches(0.3)), Inches(3.0), Inches(3.1), Inches(0.5), title, 17, WHITE, True)
    add_text(slide, Emu(x + Inches(0.3)), Inches(3.5), Inches(3.1), Inches(0.3), timing, 12, LIGHT_GRAY)
    add_accent_line(slide, Emu(x + Inches(0.3)), Inches(3.85), Inches(1.2), accent, Pt(2))
    add_bullet_list(slide, Emu(x + Inches(0.3)), Inches(4.1), Inches(3.1), Inches(2.8), items, 13, LIGHT_GRAY, (accent.red, accent.green, accent.blue) if hasattr(accent, 'red') else (0x14, 0xF1, 0x95))

# Arrow between phases
for i in range(2):
    x = Inches(4.65 + i * 4.1)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(4.5), Inches(0.3), Inches(0.3))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MED_GRAY
    arrow.line.fill.background()


# ═══════════════════════════════════════════════════════
# SLIDE 10: COMPETITION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "COMPETITION", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "No One Has Built the SSTS Implementation Layer", 38, WHITE, True)

# Competition table
headers = ["", "SSTS Compliant", "Open Source", "Solana Native", "SDK + Dashboard"]
comp_data = [
    ("TokenForge", "\u2713", "\u2713", "\u2713", "\u2713"),
    ("Securitize", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Polymath (EVM)", "\u2717", "\u2717", "\u2717", "\u2713"),
    ("Manual Anchor", "\u2713*", "\u2713*", "\u2713", "\u2717"),
    ("Hackathon Projects", "\u2717", "\u2713", "\u2713", "\u2717"),
]

table = slide.shapes.add_table(len(comp_data) + 1, len(headers), Inches(0.8), Inches(2.4), Inches(11.5), Inches(3.5)).table

for i, header in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = header
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(13)
        p.font.color.rgb = SOL_GREEN
        p.font.bold = True
        p.font.name = "Calibri"
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(0x12, 0x20, 0x30)

for row_idx, row_data in enumerate(comp_data):
    for col_idx, val in enumerate(row_data):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if row_idx == 0:
                p.font.color.rgb = SOL_GREEN
                p.font.bold = True
            elif val == "\u2713":
                p.font.color.rgb = SOL_GREEN
            elif val == "\u2717":
                p.font.color.rgb = RED_ACCENT
            else:
                p.font.color.rgb = LIGHT_GRAY
        if row_idx == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x10, 0x15, 0x2A)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_CARD if row_idx % 2 == 0 else RGBColor(0x12, 0x1A, 0x38)

add_text(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
         "* Manual Anchor = 2\u20133 months of engineering, $50K+ in dev time. TokenForge: one day, npm install.", 12, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 11: GO-TO-MARKET
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "GO-TO-MARKET", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Developer-First, Ecosystem-Driven", 38, WHITE, True)

gtm_phases = [
    ("Week 1 \u2014 Launch", [
        "Open-source release on GitHub",
        "Post to Solana Discord (#developer-tools, #defi, #rwas)",
        "Tweet thread tagging @SolanaFndn, @Halborn, @Securitize",
        "Submit Solana Foundation $250K grant application"
    ], SOL_GREEN),
    ("Month 1 \u2014 Warm Leads", [
        "Direct email to SSTS co-authors",
        "Cold outreach to C4 accelerator RWA cos",
        "Publish DeFi integration guides",
        "Kamino / Drift / Orca integration docs"
    ], SOL_PURPLE),
    ("Month 3 \u2014 Distribution", [
        "Submit to Colosseum hackathon",
        "Partnership with Squads (multisig)",
        "Breakpoint 2026 talk submission",
        "SSTS issuer case studies"
    ], ACCENT_BLUE),
]

for i, (title, items, accent) in enumerate(gtm_phases):
    x = Inches(0.8 + i * 4.1)
    add_rect(slide, x, Inches(2.5), Inches(3.7), Inches(4.2), DARK_CARD, accent, Pt(1))
    add_text(slide, Emu(x + Inches(0.3)), Inches(2.7), Inches(3.1), Inches(0.5), title, 17, accent, True)
    add_accent_line(slide, Emu(x + Inches(0.3)), Inches(3.2), Inches(1.2), accent, Pt(2))
    add_bullet_list(slide, Emu(x + Inches(0.3)), Inches(3.5), Inches(3.1), Inches(3.0), items, 14, LIGHT_GRAY, (accent.red, accent.green, accent.blue) if hasattr(accent, 'red') else (0x14, 0xF1, 0x95))

# Growth levers
add_rect(slide, Inches(0.8), Inches(6.9), Inches(11.7), Inches(0.45), DARK_CARD)
add_text(slide, Inches(1.2), Inches(6.95), Inches(11), Inches(0.4),
         "Growth lever: Every new SSTS token is a distribution event \u2014 issuers tell other issuers. DeFi integrations create pull demand.", 12, SOL_GREEN)


# ═══════════════════════════════════════════════════════
# SLIDE 12: REVENUE PROJECTIONS
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "REVENUE", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Path to $960K ARR in Year 1", 38, WHITE, True)

# Revenue milestones as bar chart
milestones = [
    ("Month 1", "$0", 0.01, "Grant-funded", SOL_GREEN),
    ("Month 3", "$5K MRR", 5, "3\u20135 Starter customers", SOL_PURPLE),
    ("Month 6", "$25K MRR", 25, "$300K ARR, mix of tiers", ACCENT_BLUE),
    ("Year 1", "$80K MRR", 80, "$960K ARR", SOL_GREEN),
]

for i, (period, amount, pct, detail, accent) in enumerate(milestones):
    x = Inches(1.5 + i * 2.9)
    bar_height = max(Inches(0.2), Inches(pct * 0.035))
    bar_y = Inches(5.5) - bar_height

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, bar_y, Inches(2.0), bar_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.adjustments[0] = 0.05

    add_text(slide, Emu(x - Inches(0.1)), Inches(5.7), Inches(2.5), Inches(0.3), period, 13, LIGHT_GRAY, True, PP_ALIGN.CENTER)
    add_text(slide, Emu(x - Inches(0.1)), bar_y - Inches(0.5), Inches(2.5), Inches(0.4), amount, 16, accent, True, PP_ALIGN.CENTER)
    add_text(slide, Emu(x - Inches(0.1)), bar_y - Inches(0.9), Inches(2.5), Inches(0.4), detail, 11, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Retention callout
add_rect(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), DARK_CARD)
add_text(slide, Inches(1.2), Inches(6.4), Inches(11), Inches(0.7),
         "Retention moat: Issuers are sticky once deployed \u2014 switching compliance infrastructure post-issuance is extremely painful. Upgrade trigger: Self-hosted issuers hit operational friction (Merkle complexity, key management) and upgrade to managed.", 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 13: RISK & VALIDATION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "RISK & VALIDATION", 12, ORANGE_ACCENT, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), ORANGE_ACCENT, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Our Riskiest Assumption \u2014 And How We\u2019ll Test It", 38, WHITE, True)

# Risk card
add_rect(slide, Inches(0.8), Inches(2.5), Inches(5.8), Inches(4.5), DARK_CARD, ORANGE_ACCENT, Pt(1))
add_text(slide, Inches(1.2), Inches(2.7), Inches(5), Inches(0.4), "THE RISK", 13, ORANGE_ACCENT, True)
add_text(slide, Inches(1.2), Inches(3.2), Inches(5), Inches(1.2),
         "Issuers will voluntarily adopt SSTS as the standard rather than rolling their own.",
         18, WHITE, True)
add_multiline(slide, Inches(1.2), Inches(4.5), Inches(5), Inches(2.5), [
    "Franklin Templeton and Securitize co-authored",
    "the spec, but institutional inertia is real \u2014",
    "they may use proprietary tooling anyway."
], 15, LIGHT_GRAY)

# Validation card
add_rect(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(4.5), DARK_CARD, SOL_GREEN, Pt(1))
add_text(slide, Inches(7.4), Inches(2.7), Inches(4.8), Inches(0.4), "VALIDATION SIGNAL", 13, SOL_GREEN, True)
add_multiline(slide, Inches(7.4), Inches(3.2), Inches(4.8), Inches(3.5), [
    "1. Email Halborn: Have any teams",
    "   requested an SDK or asked for",
    "   implementation help?",
    "",
    "   If yes \u2192 ship immediately",
    "",
    "2. Confirm Solana Foundation $250K",
    "   RFP is open and accepting",
    "   applications for SSTS tooling",
    "",
    "3. Post SSTS Anchor program on",
    "   devnet in 2 weeks \u2014 validate",
    "   the core technical claim"
], 14, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 14: THE ASK
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "THE ASK", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
         "Shipping the Canonical SSTS Implementation", 38, WHITE, True)

add_text(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.7),
         "Solana Foundation $250K RFP Grant", 28, SOL_GREEN, True)

# Use of funds breakdown
funds = [
    ("Engineering", "60%", "SSTS + FAMP programs, SDK, dashboard", Inches(4.8), SOL_GREEN),
    ("Security Audit", "20%", "Sec3 / Neodyme / OtterSec ($20\u201340K)", Inches(1.6), SOL_PURPLE),
    ("Infrastructure", "10%", "Hosting, RPC, DB ($260/mo)", Inches(0.8), ACCENT_BLUE),
    ("GTM + Community", "10%", "Events, content, partnerships", Inches(0.8), ORANGE_ACCENT),
]

y = Inches(3.3)
for label, pct, detail, bar_width, accent in funds:
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.4), f"{pct}  {label}", 16, WHITE, True)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Emu(y + Inches(0.4)), bar_width, Inches(0.35), )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.adjustments[0] = 0.15
    add_text(slide, Inches(0.8), Emu(y + Inches(0.85)), Inches(4), Inches(0.3), detail, 12, LIGHT_GRAY)
    y = Emu(y + Inches(1.2))

# First milestone
add_rect(slide, Inches(7.0), Inches(3.0), Inches(5.5), Inches(4.0), DARK_CARD, SOL_GREEN, Pt(1))
add_text(slide, Inches(7.4), Inches(3.2), Inches(4.8), Inches(0.4), "FIRST MILESTONE (2 WEEKS)", 13, SOL_GREEN, True)
add_multiline(slide, Inches(7.4), Inches(3.8), Inches(4.8), Inches(3.0), [
    "\u2022 Ship SSTS Anchor program on devnet",
    "",
    "\u2022 Working create_token + mint_to +",
    "  transfer hook",
    "",
    "\u2022 No dashboard, no SDK yet \u2014 just a",
    "  working on-chain program with",
    "  Anchor test suite",
    "",
    "\u2022 Post GitHub repo publicly on Day 14",
    "",
    "\u2022 Validates core technical claim,",
    "  establishes first-mover presence"
], 14, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 15: VISION
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Large gradient circles
shape1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1), Inches(-1), Inches(6), Inches(6))
shape1.fill.solid()
shape1.fill.fore_color.rgb = SOL_PURPLE
shape1.line.fill.background()
sp1 = shape1._element.spPr
sf1 = sp1.find(qn("a:solidFill"))
if sf1 is not None:
    sc1 = sf1.find(qn("a:srgbClr"))
    if sc1 is not None:
        a1 = etree.SubElement(sc1, qn("a:alpha"))
        a1.set("val", "8000")

shape2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8), Inches(3), Inches(7), Inches(7))
shape2.fill.solid()
shape2.fill.fore_color.rgb = SOL_GREEN
shape2.line.fill.background()
sp2 = shape2._element.spPr
sf2 = sp2.find(qn("a:solidFill"))
if sf2 is not None:
    sc2 = sf2.find(qn("a:srgbClr"))
    if sc2 is not None:
        a2 = etree.SubElement(sc2, qn("a:alpha"))
        a2.set("val", "6000")

add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.4), "THE VISION", 12, SOL_GREEN, True)
add_accent_line(slide, Inches(0.8), Inches(0.85), Inches(1), SOL_GREEN, Pt(3))

add_text(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1.2),
         "The Default Compliance Layer\nfor All RWA Tokens on Solana", 42, WHITE, True)

add_multiline(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(2.5), [
    "Just as Metaplex became the standard for NFTs on Solana,",
    "TokenForge becomes the standard for security tokens.",
    "",
    "Every SSTS token issued through TokenForge is a distribution event.",
    "DeFi integrations create pull demand.",
    "Issuers are sticky once deployed.",
], 18, LIGHT_GRAY)

add_text(slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.5),
         "npm install @tokenforge/sdk", 24, SOL_GREEN, True)

add_text(slide, Inches(0.8), Inches(6.7), Inches(6), Inches(0.4),
         "github.com/tokenforge  |  team@tokenforge.io", 14, MED_GRAY)


# ═══════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════
output_path = "/Users/reinhartsulilatu/Repos/tokenforge/pitch-deck/TokenForge_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
